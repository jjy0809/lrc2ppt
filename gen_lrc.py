from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from youtube_transcript_api import YouTubeTranscriptApi
from lrcup import LRCLib
import requests
from bs4 import BeautifulSoup
from urllib.parse import quote_plus
from ytmusicapi import YTMusic
import sys
import os
import re

def url_to_id(url):
    if "youtu" in url:
        if '.be/' in url: 
            return (url[16:url.find('?')], 'yt')
        else: 
            return (url[url.find('=')+1:url.find('&')], 'yt')
    elif "melon" in url or "kakao" in url: 
        return (url, 'ml')
    else: 
        return 0

def scr_mel(url):
    try:
        res = requests.get(url, headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36'})
        soup = BeautifulSoup(res.text, 'html.parser')
        lrc = soup.find('div', {'class': 'lyric', 'id': 'd_video_summary'})
        if lrc: 
            return [line.strip() for line in lrc.get_text('\n', strip=False).split('\n') if line.strip()]
        else: 
            return 0
    except: 
        return 0

def get_lrc(id):
    if id[1] == 'yt':
        try: 
            transcript = YouTubeTranscriptApi.get_transcript(id[0], languages=['ko'])
        except:
            print("가사 불러오기 실패")
            return 0
        return [entry['text'].strip() for entry in transcript]
    
    elif id[1] == 'ml': 
        return scr_mel(id[0])
    
    elif id[1] == 'lrc':
        title, artist = id[0].split('&&')
        lrc = LRCLib().search(track=title, artist=artist)
        if lrc: 
            return lrc[0]['plainLyrics'].split('\n')
        else:
            print("가사 불러오기 실패")
            return 0
    else: 
        print("가사 불러오기 실패")

def preprocess(lrcs, feat, len_con):
    pro_lrcs = lrcs

    for i in range(len(pro_lrcs)):
        if not pro_lrcs[i]:
            continue
        pro_lrcs[i] = re.sub(r'[^\w \' - ) (]', '', pro_lrcs[i], flags=re.UNICODE)
        pro_lrcs[i] = re.sub(r'_', '', pro_lrcs[i])
        if not feat: 
            pro_lrcs[i] = re.sub(r'[) (]', ' ', pro_lrcs[i])

    if len_con:
        tmp1 = True
        while tmp1:
            tmp1 = False
            tmp_lst = []
            i = 0
            while i < len(pro_lrcs):
                tmp2 = pro_lrcs[i]
                if not tmp2:
                    tmp_lst.append(tmp2)
                    i += 1
                    continue

                if i < len(pro_lrcs)-1 and pro_lrcs[i+1]:
                    com = f"{tmp2} {pro_lrcs[i+1]}"
                    if len(com) <= 20:
                        tmp_lst.append(com)
                        i += 2
                        tmp1 = True
                        continue
                tmp_lst.append(tmp2)
                i += 1
            pro_lrcs = tmp_lst

    for i in range(len(pro_lrcs)):
        if not pro_lrcs[i]:
            continue
        pro_lrcs[i] = re.sub(r'\s+', ' ', pro_lrcs[i].strip())
    
    return pro_lrcs

def line_split(line, ft_sz, width):
    txt_wid = ft_sz * 0.5
    max_wid = int((width.pt * 0.45) // txt_wid)

    if len(line) <= max_wid:
        return line

    mid = len(line) // 2
    lft = line.rfind(' ', 0, mid)
    rht = line.find(' ', mid)

    spl_pt = -1
    if lft != -1 and rht != -1:
        l = mid - lft
        r = rht - mid
        spl_pt = lft if l <= r else rht
    elif lft != -1:
        spl_pt = lft
    elif rht != -1:
        spl_pt = rht
    else:
        return line

    nxt_word_s = spl_pt + 1
    nxt_word_e = line.find(' ', nxt_word_s)
    nxt_word = line[nxt_word_s:nxt_word_e] if nxt_word_e != -1 else line[nxt_word_s:]

    if len(nxt_word) == 1 and nxt_word != 'I':
        ns = line.find(' ', nxt_word_s + 1)
        if ns != -1:
            spl_pt = ns
    return line[:spl_pt] + '\n' + line[spl_pt+1:].strip()

def gen_ppt(
    lrcs,
    ft_sz = 60,
    ft_name = 'Malgun Gothic',
    ft_col = (255,255,255),
    bg_img = None,
    ppt_name = 'lyrics.pptx',
    bold = True,
    italic = False,
    empt = False,
    pre = True,
    feat = False,
    len_con = False,
    lin_space = 1.4
):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    if pre: 
        lrcs = preprocess(lrcs = lrcs, feat = feat, len_con = len_con)

    for lrc in lrcs:
        if not empt and not lrc.strip():
            continue

        lrc = line_split(lrc, ft_sz, prs.slide_width)

        sld = prs.slides.add_slide(prs.slide_layouts[6])

        if bg_img and os.path.exists(bg_img):
            pic = sld.shapes.add_picture(
                bg_img, 0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            sld.shapes._spTree.insert(2, pic._element)

        text_box = sld.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = lrc.replace('\n', '\v')
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = lin_space
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

        font = paragraph.font
        font.name = ft_name
        font.size = Pt(ft_sz)
        font.bold = bold
        font.italic = italic
        font.color.rgb = RGBColor(*ft_col)

        text_box.top = int((prs.slide_height - text_box.height) / 2)

    prs.save(ppt_name)

if __name__ == "__main__":
    lrcs = 0

    while lrcs == 0:
        url = input("음원 링크 (유튜브, 유튜브 뮤직, 멜론) 또는 곡 제목과 아티스트 이름 (곡&&아티스트) -> ")
        if url == '0': 
            break
        
        if '&&' in url: 
            id = (url, 'lrc')
        else: 
            id = url_to_id(url)
        print(id)
        lrcs = get_lrc(id)
        print(lrcs)

        if lrcs == 0: 
            lrcs = sys.stdin.read().splitlines()

        gen_ppt(
            lrcs = lrcs,
            ft_col = (255, 255, 255),
            bg_img = r'C:\Users\happy\Desktop\학교\고등학교\3학년\가사 추출\black.jpg',
            ppt_name = r'C:\Users\happy\Desktop\학교\고등학교\3학년\가사 추출\res.pptx',
            empt = False,
            lin_space = 1.4
        )
