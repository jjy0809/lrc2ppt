# main.py
from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from youtube_transcript_api import YouTubeTranscriptApi
from lrcup import LRCLib
import requests
from bs4 import BeautifulSoup
from urllib.parse import quote_plus
from ytmusicapi import YTMusic
import re
import os
import uuid
import asyncio
from pydantic import BaseModel
from typing import Optional, List
import uvicorn

app = FastAPI()

from fastapi.middleware.cors import CORSMiddleware

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://lrc.kro.kr"],
    allow_methods=["*"],
    allow_headers=["*"]
)

# 요청 데이터 모델
class LinkRequest(BaseModel):
    link: str

class PPTSettings(BaseModel):
    lyrics: List[str]
    ft_sz: int = 60
    ft_name: str = 'Malgun Gothic'
    ft_col: tuple = (255,255,255)
    bg_img: Optional[str] = None
    bold: bool = True
    italic: bool = False
    empt: bool = False
    pre: bool = True
    feat: bool = False
    len_con: bool = False
    lin_space: float = 1.4
    artist: Optional[str] = "Unknown"  # 아티스트 필드 추가
    title: Optional[str] = "Untitled"  # 제목 필드 추가

# URL 파싱 함수
def url_to_id(url):
    if "youtu" in url:
        if '.be/' in url: 
            video_id = url.split('.be/')[1].split('?')[0]
            return (video_id, 'yt')
        else: 
            video_id = url.split('v=')[1].split('&')[0]
            return (video_id, 'yt')
    elif "melon" in url or "kakao" in url:
        return (url, 'ml')
    else:
        return (None, 'unknown')

# 멜론 가사 스크래핑 (비동기)
async def scr_mel(url):
    try:
        loop = asyncio.get_running_loop()
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36'}
        res = await loop.run_in_executor(None, lambda: requests.get(url, headers=headers, timeout=10))
        soup = BeautifulSoup(res.text, 'html.parser')
        lrc = soup.find('div', {'class': 'lyric', 'id': 'd_video_summary'})
        return [line.strip() for line in lrc.get_text('\n', strip=False).split('\n') if line.strip()] if lrc else None
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Melon 가사 추출 실패: {str(e)}")

# 가사 추출 메인 함수
async def get_lrc(id):
    try:
        if id[1] == 'yt':
            transcript = await asyncio.to_thread(YouTubeTranscriptApi.get_transcript, id[0], languages=['ko'])
            return [entry['text'].strip() for entry in transcript]
        elif id[1] == 'ml':
            return await scr_mel(id[0])
        elif id[1] == 'lrc':
            title, artist = id[0].split('&&')
            lrc = LRCLib().search(track=title, artist=artist)
            return lrc[0]['plainLyrics'].split('\n') if lrc else None
        else:
            return None
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"가사 추출 실패: {str(e)}")

# 가사 전처리 (원본 코드 그대로)
async def preprocess(lrcs, feat, len_con):
    pro_lrcs = lrcs.copy()
    
    for i in range(len(pro_lrcs)):
        if not pro_lrcs[i]:
            continue
        # 특수문자 제거
        pro_lrcs[i] = re.sub(r'[^\w \' - ) (]', '', pro_lrcs[i], flags=re.UNICODE)
        pro_lrcs[i] = re.sub(r'_', '', pro_lrcs[i])
        if not feat:
            pro_lrcs[i] = re.sub(r'[) (]', ' ', pro_lrcs[i])

    if len_con:
        tmp1 = True
        while tmp1:
            tmp1 = False
            tmp_lst = []
            i = 0
            while i < len(pro_lrcs):
                tmp2 = pro_lrcs[i]
                if not tmp2:
                    tmp_lst.append(tmp2)
                    i += 1
                    continue

                if i < len(pro_lrcs)-1 and pro_lrcs[i+1]:
                    com = f"{tmp2} {pro_lrcs[i+1]}"
                    if len(com) <= 20:
                        tmp_lst.append(com)
                        i += 2
                        tmp1 = True
                        continue
                tmp_lst.append(tmp2)
                i += 1
            pro_lrcs = tmp_lst

    for i in range(len(pro_lrcs)):
        if pro_lrcs[i]:
            pro_lrcs[i] = re.sub(r'\s+', ' ', pro_lrcs[i].strip())

    return pro_lrcs

# 줄바꿈 처리 (원본 코드 그대로)
async def line_split(line, ft_sz, width):
    txt_wid = ft_sz * 0.5
    max_wid = int((width.pt * 0.45) // txt_wid)

    if len(line) <= max_wid:
        return line

    mid = len(line) // 2
    lft = line.rfind(' ', 0, mid)
    rht = line.find(' ', mid)

    spl_pt = -1
    if lft != -1 and rht != -1:
        l = mid - lft
        r = rht - mid
        spl_pt = lft if l <= r else rht
    elif lft != -1:
        spl_pt = lft
    elif rht != -1:
        spl_pt = rht
    else:
        return line

    nxt_word_s = spl_pt + 1
    nxt_word_e = line.find(' ', nxt_word_s)
    nxt_word = line[nxt_word_s:nxt_word_e] if nxt_word_e != -1 else line[nxt_word_s:]

    if len(nxt_word) == 1 and nxt_word != 'I':
        ns = line.find(' ', nxt_word_s + 1)
        if ns != -1:
            spl_pt = ns

    return line[:spl_pt] + '\n' + line[spl_pt+1:].strip()

# PPT 생성 함수 (원본 코드 수정 버전)
def gen_ppt_sync(settings: PPTSettings):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    if settings.pre:
        loop = asyncio.new_event_loop()
        lrcs = loop.run_until_complete(preprocess(settings.lyrics, settings.feat, settings.len_con))
    else:
        lrcs = settings.lyrics

    for lrc_line in lrcs:
        if not settings.empt and not lrc_line.strip():
            continue

        line = loop.run_until_complete(line_split(
            lrc_line, 
            settings.ft_sz, 
            prs.slide_width
        ))

        sld = prs.slides.add_slide(prs.slide_layouts[6])

        if settings.bg_img and os.path.exists(settings.bg_img):
            pic = sld.shapes.add_picture(
                settings.bg_img, 0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            sld.shapes._spTree.insert(2, pic._element)

        text_box = sld.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        paragraph = text_frame.paragraphs[0]
        paragraph.text = line.replace('\n', '\v')
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = settings.lin_space
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

        font = paragraph.font
        font.name = settings.ft_name
        font.size = Pt(settings.ft_sz)
        font.bold = settings.bold
        font.italic = settings.italic
        font.color.rgb = RGBColor(*settings.ft_col)

        text_box.top = int((prs.slide_height - text_box.height) / 2)

    # 파일명 생성: 아티스트-곡제목.pptx
    safe_artist = re.sub(r'[\\/*?:"<>|]', "", settings.artist)
    safe_title = re.sub(r'[\\/*?:"<>|]', "", settings.title)
    filename = f"{safe_artist}-{safe_title}.pptx"
    prs.save(filename)
    
    return filename

# API 엔드포인트
@app.post("/api/1")
async def api_1(request: LinkRequest):
    try:
        if '&&' in request.link:
            id = (request.link, 'lrc')
        else:
            id = url_to_id(request.link)
            if id[1] == 'unknown':
                raise ValueError("지원하지 않는 플랫폼")

        lrcs = await get_lrc(id)
        if not lrcs:
            raise HTTPException(status_code=404, detail="가사를 찾을 수 없음")

        processed = await preprocess(lrcs, feat=False, len_con=True)
        return {"lyrics": processed}
    
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/2")
async def api_2(settings: PPTSettings):
    try:
        filename = await asyncio.to_thread(gen_ppt_sync, settings)
        return FileResponse(
            filename,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=filename  # 다운로드 시 원본 파일명 사용
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(filename):
            os.remove(filename)  # 임시 파일 삭제

if __name__ == "__main__":
    uvicorn.run("lrc_api:app", host="0.0.0.0", port=8000, reload=True)
